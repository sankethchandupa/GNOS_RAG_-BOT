from langchain_community.document_loaders import PyPDFLoader
from langchain_text_splitters import RecursiveCharacterTextSplitter
from sentence_transformers import SentenceTransformer
from duckduckgo_search import DDGS
from groq import Groq
from langdetect import detect
from deep_translator import GoogleTranslator
import faiss
import numpy as np
import os
import pickle
import re

from dotenv import load_dotenv
load_dotenv()

from huggingface_hub import login
login(token=os.getenv("HF_TOKEN"))

from docx import Document
from pptx import Presentation


client = Groq(api_key=os.getenv("GROQ_API_KEY"))


print("Loading models...")
embedding_model = SentenceTransformer("intfloat/multilingual-e5-small")


FAISS_FILE = "faiss_index.index"
DOCS_FILE  = "docs.pkl"


SYSTEM_PROMPT = """You are a highly accurate, helpful assistant.

IMPORTANT RULES:
1. The context provided may contain raw PDF text including application form fields
   such as "Name:", "Gender:", "Date:", "As appear in NIC", "Line 1", "City" etc.
   COMPLETELY IGNORE any form fields, blank fields, labels, or form structure.
2. If the context does not contain a clear, useful answer to the question,
   DO NOT make up an answer from the context. Instead answer from your own knowledge.
3. Always write your answer as clean, clear, well-structured prose or numbered points.
4. NEVER copy raw form field names, blank fields, or form structure into your answer.
5. Answer in {answer_language} language.
6. Be complete, accurate, and easy to understand.
"""


# =========================
def contains_sinhala(text: str) -> bool:
    return any('\u0D80' <= ch <= '\u0DFF' for ch in text)

def contains_tamil(text: str) -> bool:
    return any('\u0B80' <= ch <= '\u0BFF' for ch in text)


def detect_language(text: str) -> str:
    """
    Detect language of text.
    First checks Unicode script ranges (reliable for Sinhala/Tamil).
    Falls back to langdetect for other languages.
    """
    
    if contains_sinhala(text):
        return "Sinhala"
    if contains_tamil(text):
        return "Tamil"

    
    try:
        lang = detect(text)
        if lang == "si":
            return "Sinhala"
        elif lang == "ta":
            return "Tamil"
        else:
            return "English"
    except:
        return "English"

def detect_preferred_language(question: str) -> str:
    """
    Detect what language the user wants the answer in.
    Priority:
    1. Explicit keyword in question (e.g. "give sinhala answer")
    2. Unicode script detection (Sinhala/Tamil characters in question)
    3. langdetect fallback
    """
    q_lower = question.lower()

    
    if any(kw in q_lower for kw in ["sinhala", "සිංහල", "in si", "sinhala language", "sinhalen"]):
        return "Sinhala"
    elif any(kw in q_lower for kw in ["tamil", "தமிழ்", "in ta", "tamil language"]):
        return "Tamil"
    elif any(kw in q_lower for kw in ["english", "in english", "english language"]):
        return "English"

    
    return detect_language(question)


def is_context_useful(context: str) -> bool:
    """
    Returns True  if context looks like real informational text.
    Returns False if context looks like a form / label dump.
    If more than 40% of lines look like form fields → not useful.
    """
    if not context.strip():
        return False

    lines = [l.strip() for l in context.split('\n') if l.strip()]
    if not lines:
        return False

    form_like = 0
    for line in lines:
        if len(line) < 50 and line.endswith(':'):
            form_like += 1
            continue
        ascii_ratio = sum(1 for c in line if ord(c) < 128) / max(len(line), 1)
        if ascii_ratio > 0.9 and len(line) < 40:
            form_like += 1
            continue
        if re.search(
            r'\b(gender|male|female|married|unmarried|divorced|widow|'
            r'occupation|purpose|passport|visa|scholarship|employment|'
            r'line \d|postal|zip code|state\/province|yes\s+no|if yes|'
            r'as appear|issued date|reference no|overseas)\b',
            line, re.IGNORECASE
        ):
            form_like += 1

    return (form_like / max(len(lines), 1)) < 0.4


def translate_to_english(text: str) -> str:
    try:
        return GoogleTranslator(source='auto', target='en').translate(text)
    except:
        return text

def translate_answer(text: str, target_lang: str) -> str:
    """
    Translate clean English answer to Sinhala or Tamil.
    Groups lines into 800-char chunks so Google Translate gets
    full sentence context — prevents letter-by-letter breakdown.
    """
    try:
        if target_lang == "Sinhala":
            lang_code = 'si'
        elif target_lang == "Tamil":
            lang_code = 'ta'
        else:
            return text

        lines = text.split('\n')
        chunks = []
        current_chunk_lines = []
        current_len = 0
        MAX_CHUNK = 800

        for line in lines:
            line_len = len(line)
            if current_len + line_len + 1 > MAX_CHUNK and current_chunk_lines:
                chunks.append('\n'.join(current_chunk_lines))
                current_chunk_lines = []
                current_len = 0
            current_chunk_lines.append(line)
            current_len += line_len + 1

        if current_chunk_lines:
            chunks.append('\n'.join(current_chunk_lines))

        translated_chunks = []
        for chunk in chunks:
            chunk = chunk.strip()
            if not chunk:
                translated_chunks.append('')
                continue
            try:
                translated = GoogleTranslator(source='en', target=lang_code).translate(chunk)
                translated_chunks.append(translated if translated else chunk)
            except:
                translated_chunks.append(chunk)

        return '\n'.join(translated_chunks)

    except:
        return text


def smart_translate(raw_answer: str, answer_language: str) -> str:
    if answer_language == "English":
        return raw_answer
    elif answer_language == "Sinhala":
        
        return translate_answer(raw_answer, "Sinhala")
    elif answer_language == "Tamil":
        if contains_tamil(raw_answer):
            return raw_answer
        return translate_answer(raw_answer, "Tamil")
    return raw_answer


def search_web(query: str) -> str:
    try:
        with DDGS() as ddgs:
            results = ddgs.text(query, max_results=5)
            return "\n\n".join([r["body"] for r in results if r.get("body")])
    except:
        return ""


chat_history = []
MAX_MEMORY = 5

def update_chat_history(q, a):
    chat_history.append((q, a))
    if len(chat_history) > MAX_MEMORY:
        chat_history.pop(0)

def get_chat_history():
    return "\n".join([f"User: {q}\nAssistant: {a}" for q, a in chat_history])


def load_pdf(path):
    try:
        loader = PyPDFLoader(path)
        docs = loader.load()
        return [d.page_content for d in docs]
    except Exception as e:
        print(f" Error loading PDF {path}: {e}")
        return []

def load_txt(path):
    with open(path, "r", encoding="utf-8") as f:
        return [f.read()]

def load_docx(path):
    doc = Document(path)
    return ["\n".join([p.text for p in doc.paragraphs])]

def load_pptx(path):
    prs = Presentation(path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return ["\n".join(text)]

def load_file(path):
    ext = path.lower().split(".")[-1]
    if ext == "pdf":    return load_pdf(path)
    elif ext == "txt":  return load_txt(path)
    elif ext == "docx": return load_docx(path)
    elif ext == "pptx": return load_pptx(path)
    else:               return []


def build_index(paths):
    splitter = RecursiveCharacterTextSplitter(chunk_size=800, chunk_overlap=150)
    docs = []

    for path in paths:
        contents = load_file(path)
        for content in contents:
            translated_content = translate_to_english(content)
            chunks = splitter.split_text(translated_content)
            for c in chunks:
                docs.append({"text": c, "source": path})

    texts = [d["text"] for d in docs]
    embeddings = embedding_model.encode(texts, batch_size=32, show_progress_bar=True)
    emb_array = np.array(embeddings).astype("float32")

    index = faiss.IndexFlatL2(emb_array.shape[1])
    index.add(emb_array)
    return index, docs


pdf_folder = "pdfs"
paths = []

for file_name in os.listdir(pdf_folder):
    file_path = os.path.join(pdf_folder, file_name)
    if os.path.isfile(file_path) and file_name.lower().endswith(".pdf"):
        print("Loading file:", file_name)
        paths.append(file_path)


if os.path.exists(FAISS_FILE) and os.path.exists(DOCS_FILE):
    print("Loading saved FAISS index...")
    index = faiss.read_index(FAISS_FILE)
    with open(DOCS_FILE, "rb") as f:
        docs = pickle.load(f)
    print(f"Loaded successfully. Vectors: {index.ntotal}")
else:
    print("  Index not found. Building index FIRST TIME...")
    index, docs = build_index(paths)
    faiss.write_index(index, FAISS_FILE)
    with open(DOCS_FILE, "wb") as f:
        pickle.dump(docs, f)
    print(" Index built and saved successfully")


def rag_pipeline(question: str) -> str:
    """
    Main entry point called by api.py.
    Detects language, retrieves context, generates and translates answer.
    """
    answer_language = detect_preferred_language(question)
    print(f" Detected language: {answer_language}")

    english_question = translate_to_english(question)

    
    q_emb = embedding_model.encode([english_question])
    q_emb = np.array(q_emb).astype("float32")
    distances, indices_result = index.search(q_emb, 5)

    best_score = distances[0][0]
    top_chunks = [docs[i] for i in indices_result[0][:3]]
    pdf_context = "\n\n".join([c["text"] for c in top_chunks])

    
    pdf_is_relevant = best_score < 1.0
    pdf_is_clean    = is_context_useful(pdf_context)

    if pdf_is_relevant and pdf_is_clean:
        final_context = pdf_context
        print(" Using PDF context")
    else:
        print(" PDF context not useful. Searching web...")
        web_data = search_web(english_question)
        if web_data.strip():
            final_context = web_data
            print(" Using web context")
        else:
            final_context = ""
            print(" Using LLM knowledge only")

    
    llm_answer_language = "English" if answer_language == "Sinhala" else answer_language
    system_prompt = SYSTEM_PROMPT.format(answer_language=llm_answer_language)

    if final_context.strip():
        user_content = f"Context:\n{final_context}\n\nQuestion: {english_question}"
    else:
        user_content = f"Question: {english_question}"

    response = client.chat.completions.create(
        model="llama-3.1-8b-instant",
        messages=[
            {"role": "system", "content": system_prompt},
            {"role": "user",   "content": user_content}
        ],
        max_tokens=1024,
        temperature=0.2
    )

    raw_answer = response.choices[0].message.content

    
    final_answer = smart_translate(raw_answer, answer_language)

    update_chat_history(question, final_answer)
    return final_answer



if __name__ == "__main__":

    while True:
        question = input("\nAsk (or type exit / clear): ").strip()

        if not question:
            continue
        if question.lower() == "exit":
            break
        if question.lower() == "clear":
            chat_history.clear()
            print(" Chat history cleared!")
            continue

        answer = rag_pipeline(question)
        print(f"\n Answer:\n{answer}")
